"""
PresentiQ — AI-powered presentation generator.

Pipeline:
  1. Two-stage outline generation (document analysis → outline)
  2. HTML/CSS slide rendering via LLM + Playwright screenshot
  3. Style-anchored batch generation for visual consistency
  4. Caching for outlines and images
  5. Smart error handling with retries and fallbacks
"""

import asyncio
import os
import json
import logging
from datetime import datetime
from typing import List, Dict, Optional

from .outline_generator import OutlineGenerator
from .claude_client import ClaudeClient, OpenaiClient, UnifiedAIClient
from .cache_manager import CacheManager
from .batch_generator import BatchImageGenerator
from .error_handler import SmartErrorHandler
from .template_loader import get_template_presets
from .slide_generator_official import SlideGenerator, ImageGenerationTool

__version__ = "2.0.0"

logger = logging.getLogger(__name__)


class PPTGenerator:
    """
    PresentiQ main class — transforms text into professional slide decks.

    Core features:
    - Two-stage outline generation (document analysis + outline generation)
    - Style-anchored batch image generation
    - Caching (outline + images)
    - Smart error handling and fallback
    - 46 template presets
    """

    def __init__(
        self,
        api_key: str = None,
        provider: str = "Claude",
        base_url: str = None,
        enable_cache: bool = True,
        cache_dir: str = ".cache/ppt_generator"
    ):
        """
        Initialize PresentiQ generator

        Args:
            api_key: API key, if None read from environment
            provider: AI provider (Claude/Openai/other)
            base_url: Custom API base URL
            enable_cache: Whether to enable cache
            cache_dir: Cache directory
        """
        if provider == "Claude":
            self.llm_client = ClaudeClient(api_key)
        elif provider == "Openai":
            self.llm_client = OpenaiClient(api_key=api_key, base_url=base_url)
        else:
            self.llm_client = UnifiedAIClient(api_key=api_key, base_url=base_url)

        self.outline_generator = OutlineGenerator(self.llm_client)
        self.image_tool = ImageGenerationTool(llm_client=self.llm_client)
        self.slide_generator = SlideGenerator(self.image_tool)

        # Enhanced components
        self.cache_manager = CacheManager(cache_dir) if enable_cache else None
        self.batch_generator = BatchImageGenerator(self.image_tool, self.cache_manager)
        self.error_handler = SmartErrorHandler()

        self.enable_cache = enable_cache

    def generate_ppt(
        self,
        reference_text: str,
        style_requirements: str,
        output_dir: str = "output",
        model: str = "deepseek-chat",
        audience_profile: Dict = None,
        brand_guidelines: Dict = None,
        brand_references: List[str] = None,
        use_cache: bool = True,
        template_preset: str = None,
        _extra_slides: List[Dict] = None,
        persona_context: Dict = None,
    ) -> dict:
        """
        Synchronously generate PPT

        Args:
            reference_text: Reference text
            style_requirements: Style requirements
            output_dir: Output directory
            model: Model to use
            audience_profile: Target audience info
            brand_guidelines: Brand guidelines
            brand_references: Brand reference image paths (max 14)
            use_cache: Whether to use cache
            template_preset: Template preset name
            _extra_slides: Optional list of editable slide specs appended
                after the generated image slides (e.g. team members, thank you).
            persona_context: Persona/audience context from PersonaEngine (optional)

        Returns:
            dict: Generation result with file paths and info
        """
        return asyncio.run(self.generate_ppt_async(
            reference_text, style_requirements, output_dir, model,
            audience_profile, brand_guidelines, brand_references, use_cache,
            template_preset=template_preset,
            _extra_slides=_extra_slides,
            persona_context=persona_context,
        ))

    async def generate_ppt_async(
        self,
        reference_text: str,
        style_requirements: str,
        output_dir: str = "output",
        model: str = "deepseek-chat",
        audience_profile: Dict = None,
        brand_guidelines: Dict = None,
        brand_references: List[str] = None,
        use_cache: bool = True,
        max_concurrent: int = 4,
        template_preset: str = None,
        _extra_slides: List[Dict] = None,
        persona_context: Dict = None,
    ) -> dict:
        """
        Asynchronously generate PPT (recommended)

        Args:
            reference_text: Reference text
            style_requirements: Style requirements
            output_dir: Output directory
            model: Model to use
            audience_profile: Target audience info
            brand_guidelines: Brand guidelines
            brand_references: Brand reference image paths
            use_cache: Whether to use cache
            max_concurrent: Max concurrency
            template_preset: Template preset name
            _extra_slides: Optional editable slides appended after image slides
            persona_context: Persona/audience context from PersonaEngine (optional)

        Returns:
            dict: Generation result with file paths and details
        """
        generation_info = {
            "mode": "enhanced",
            "two_stage": True,
            "cache_used": False,
            "style_anchored": True,
            "template_preset": template_preset,
            "persona_context": persona_context,
        }

        # ===== Phase 1: Check cache =====
        outline_result = None
        if use_cache and self.cache_manager:
            logger.info("Checking outline cache...")
            outline_result = self.cache_manager.get_cached_outline(
                reference_text, style_requirements, model
            )
            if outline_result:
                logger.info(f"Outline cache hit, {len(outline_result.get('slides', []))} slides")
                generation_info["cache_used"] = True

        # ===== Phase 2: Two-stage outline generation =====
        if not outline_result:
            logger.info("Starting two-stage outline generation...")
            try:
                outline_result = self.outline_generator.generate_outline_two_stage(
                    reference_text=reference_text,
                    style_requirements=style_requirements,
                    audience_profile=audience_profile,
                    brand_guidelines=brand_guidelines,
                    model=model,
                    template_preset=template_preset,
                    persona_context=persona_context,
                )
                logger.info(f"Two-stage outline generation complete, {len(outline_result.get('slides', []))} slides")

                # Cache outline
                if use_cache and self.cache_manager:
                    self.cache_manager.cache_outline(
                        reference_text, style_requirements, model, outline_result
                    )
                    logger.info("Outline cached")

            except Exception as e:
                logger.warning(f"Two-stage generation failed, falling back to single-stage: {e}")
                outline_result = self.outline_generator.generate_outline(
                    reference_text, style_requirements, model=model, template_preset=template_preset
                )
                generation_info["two_stage"] = False

        # ===== Phase 3: Batch image generation (with style anchoring) =====
        logger.info("Starting batch image generation (style anchoring mode)...")

        # Get template preset style hints
        style_hints = None
        template_presets = get_template_presets()
        if template_preset and template_preset in template_presets:
            preset = template_presets[template_preset]
            style_hints = preset.get('style_hints')
            if style_hints:
                logger.info(f"Using preset style: {preset.get('name', template_preset)} - {style_hints.get('background', '')}")

        try:
            slides = outline_result.get('slides', [])
            batch_results = await self.batch_generator.generate_with_style_consistency(
                slides=slides,
                outline_result=outline_result,
                style_requirements=style_requirements,
                output_dir=output_dir,
                brand_references=brand_references,
                max_concurrent=max_concurrent,
                style_hints=style_hints,
                persona_context=persona_context,
            )

            # Convert to standard format
            slides_data = []
            for result in batch_results:
                slides_data.append({
                    'success': result.success,
                    'file_path': result.file_path,
                    'error': result.error,
                    'title': slides[result.slide_index].get('title', '') if result.slide_index < len(slides) else '',
                    'from_cache': result.from_cache,
                    'retries': result.retries
                })

        except Exception as e:
            logger.warning(f"Batch generation failed, falling back to standard mode: {e}")
            generation_info["style_anchored"] = False
            slides_data = await self._generate_ppt_slides_fallback(
                outline_result, style_requirements, output_dir, max_concurrent,
                style_hints=style_hints  # Pass style_hints to fallback mode
            )

        # ===== Phase 4: Save results =====
        if _extra_slides:
            outline_result["_extra_slides"] = _extra_slides
        result = self._save_ppt(slides_data, outline_result, output_dir)
        result["generation_info"] = generation_info

        # Count cache hits
        cache_hits = sum(1 for s in slides_data if s.get('from_cache', False))
        if cache_hits > 0:
            result["cache_hits"] = cache_hits
            logger.info(f"Image cache hits: {cache_hits}/{len(slides_data)}")

        return result

    async def _generate_ppt_slides_fallback(
        self,
        outline_result: dict,
        style_requirements: str,
        output_dir: str,
        max_concurrent: int = 4,
        timeout_per_slide: float = 120.0,
        style_hints: dict = None
    ) -> List[dict]:
        """
        Fallback mode: Concurrent PPT slide generation (no style anchoring)

        Args:
            outline_result: Outline result dict
            style_requirements: Style requirements
            output_dir: Output directory
            max_concurrent: Max concurrency
            timeout_per_slide: Timeout per slide (seconds)
            style_hints: Template preset style hints (optional)

        Returns:
            List[dict]: Slide data list
        """
        slides = outline_result['slides']
        total_slides = len(slides)

        logger.info(f"Fallback mode: Concurrent PPT generation, {total_slides} slides, concurrency: {max_concurrent}")
        if style_hints:
            logger.info(f"Using preset style hints: {style_hints.get('background', '')[:30]}...")

        semaphore = asyncio.Semaphore(max_concurrent)

        async def generate_single_slide(i: int, slide_info: dict) -> tuple:
            async with semaphore:
                try:
                    result = await asyncio.wait_for(
                        self.slide_generator.generate_slide_as_image(
                            slide_info, i, outline_result, style_requirements, output_dir,
                            style_hints=style_hints  # Pass style_hints
                        ),
                        timeout=timeout_per_slide
                    )

                    if result.get("success"):
                        logger.info(f"Slide {i+1}/{total_slides} generated successfully - {slide_info.get('title', 'Untitled')}")
                        return (i, {'success': True, 'file_path': result.get('file_path'), 'title': slide_info.get('title', '')})
                    else:
                        error_msg = result.get('error', 'Unknown error')
                        logger.warning(f"Slide {i+1}/{total_slides} generation failed: {error_msg}")
                        return (i, {'success': False, 'error': error_msg, 'title': slide_info.get('title', '')})

                except asyncio.TimeoutError:
                    logger.error(f"Slide {i+1}/{total_slides} generation timeout")
                    return (i, {'success': False, 'error': 'Generation timeout', 'title': slide_info.get('title', '')})

                except Exception as e:
                    logger.exception(f"Slide {i+1}/{total_slides} generation exception: {str(e)}")
                    return (i, {'success': False, 'error': str(e), 'title': slide_info.get('title', '')})

        tasks = [generate_single_slide(i, slide) for i, slide in enumerate(slides)]
        results = await asyncio.gather(*tasks)

        sorted_results = sorted(results, key=lambda x: x[0])
        slides_data = [data for _, data in sorted_results]

        logger.info(f"All {total_slides} slides generated successfully!")
        return slides_data

    def _save_ppt(self, slides_data: list, outline: dict, output_dir: str) -> dict:
        """
        Save PPT as .pptx file

        Args:
            slides_data: Slide data list (with image paths or error info)
            outline: Outline info
            output_dir: Output directory

        Returns:
            dict: Save result info
                - pptx_file: .pptx file path
                - outline_file: Outline JSON path
                - total_slides: Total slide count
                - success_slides: Success count
                - error_slides: Failed slide list
                - timestamp: Timestamp
        """
        from pptx import Presentation
        from pptx.util import Inches

        os.makedirs(output_dir, exist_ok=True)
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")

        # Create PPT object (16:9 ratio)
        prs = Presentation()
        prs.slide_width = Inches(13.333)   # 16:9 standard width
        prs.slide_height = Inches(7.5)     # 16:9 standard height

        success_count = 0
        error_slides = []

        for i, slide_data in enumerate(slides_data):
            # Add blank slide
            blank_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(blank_layout)

            if slide_data.get('success') and slide_data.get('file_path'):
                image_path = slide_data['file_path']
                if os.path.exists(image_path):
                    # Successful slide: insert image as full-screen background
                    slide.shapes.add_picture(
                        image_path,
                        Inches(0), Inches(0),
                        width=prs.slide_width,
                        height=prs.slide_height
                    )
                    success_count += 1
                else:
                    # Image file does not exist
                    self._add_error_text_to_slide(
                        slide, prs, i + 1,
                        slide_data.get('title', ''),
                        f"Image file not found: {image_path}"
                    )
                    error_slides.append({
                        'page': i + 1,
                        'title': slide_data.get('title', ''),
                        'error': f"Image file not found: {image_path}"
                    })
            else:
                # Failed slide: add error text
                error_msg = slide_data.get('error', 'Unknown error')
                self._add_error_text_to_slide(
                    slide, prs, i + 1,
                    slide_data.get('title', ''),
                    error_msg
                )
                error_slides.append({
                    'page': i + 1,
                    'title': slide_data.get('title', ''),
                    'error': error_msg
                })

        # Append extra editable slides (team members, thank you, etc.)
        extra_slides = outline.get("_extra_slides", [])
        for extra in extra_slides:
            self._add_editable_slide(prs, extra)

        # Save PPT file
        pptx_filename = f"presentation_{timestamp}.pptx"
        pptx_path = os.path.join(output_dir, pptx_filename)
        prs.save(pptx_path)

        # Save outline info
        outline_file = os.path.join(output_dir, f"outline_{timestamp}.json")
        with open(outline_file, 'w', encoding='utf-8') as f:
            json.dump(outline, f, ensure_ascii=False, indent=2)

        total = len(slides_data) + len(extra_slides)
        logger.info(f"PPT saved: {pptx_path}, success {success_count}/{len(slides_data)} image slides + {len(extra_slides)} editable slides")

        return {
            'pptx_file': pptx_path,
            'outline_file': outline_file,
            'total_slides': total,
            'success_slides': success_count + len(extra_slides),
            'error_slides': error_slides,
            'timestamp': timestamp
        }

    def _add_error_text_to_slide(
        self,
        slide,
        prs,
        page_number: int,
        title: str,
        error_msg: str
    ):
        """
        Add error message text to slide

        Args:
            slide: Slide object
            prs: Presentation object
            page_number: Page number
            title: Slide title
            error_msg: Error message
        """
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        # Add red background shape
        shape = slide.shapes.add_shape(
            1,  # MSO_SHAPE.RECTANGLE
            Inches(0), Inches(0),
            prs.slide_width, prs.slide_height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(220, 53, 69)  # Bootstrap danger red
        shape.line.fill.background()

        # Add title text box
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2.5), Inches(11.333), Inches(1)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Slide Generation Error"
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

        # Add page info
        info_box = slide.shapes.add_textbox(
            Inches(1), Inches(3.8), Inches(11.333), Inches(0.8)
        )
        tf2 = info_box.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = f"Slide {page_number} - {title or 'Untitled'}"
        p2.font.size = Pt(28)
        p2.font.color.rgb = RGBColor(255, 255, 255)
        p2.alignment = PP_ALIGN.CENTER

        # Add error details
        detail_box = slide.shapes.add_textbox(
            Inches(1), Inches(5), Inches(11.333), Inches(1)
        )
        tf3 = detail_box.text_frame
        p3 = tf3.paragraphs[0]
        p3.text = f"Error: {error_msg}"
        p3.font.size = Pt(18)
        p3.font.color.rgb = RGBColor(255, 200, 200)
        p3.alignment = PP_ALIGN.CENTER

    def _add_editable_slide(self, prs, slide_spec: dict):
        """
        Add a fully editable text slide (not an image) to the presentation.

        Args:
            prs: Presentation object
            slide_spec: Dict with keys:
                - type: "team_members" | "thank_you" | "custom"
                - title: Slide title
                - subtitle: Optional subtitle
                - items: Optional list of text items (for team members, etc.)
                - bg_color: Background RGB tuple, e.g. (79, 195, 247)
                - text_color: Text RGB tuple, e.g. (255, 255, 255)
        """
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        bg_rgb = slide_spec.get("bg_color", (79, 195, 247))
        text_rgb = slide_spec.get("text_color", (255, 255, 255))
        accent_rgb = slide_spec.get("accent_color", (255, 255, 255))

        # Background
        bg = slide.shapes.add_shape(
            1, Inches(0), Inches(0), prs.slide_width, prs.slide_height
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(*bg_rgb)
        bg.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(0.8), Inches(11.333), Inches(1.2)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = slide_spec.get("title", "")
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*text_rgb)
        p.alignment = PP_ALIGN.CENTER

        # Subtitle
        if slide_spec.get("subtitle"):
            sub_box = slide.shapes.add_textbox(
                Inches(1), Inches(2.0), Inches(11.333), Inches(0.8)
            )
            tf2 = sub_box.text_frame
            tf2.word_wrap = True
            p2 = tf2.paragraphs[0]
            p2.text = slide_spec["subtitle"]
            p2.font.size = Pt(24)
            p2.font.color.rgb = RGBColor(*accent_rgb)
            p2.alignment = PP_ALIGN.CENTER

        # Items list (editable placeholders)
        items = slide_spec.get("items", [])
        if items:
            start_y = 3.0
            for idx, item in enumerate(items):
                item_box = slide.shapes.add_textbox(
                    Inches(2), Inches(start_y + idx * 0.7),
                    Inches(9.333), Inches(0.6)
                )
                tf_item = item_box.text_frame
                tf_item.word_wrap = True
                p_item = tf_item.paragraphs[0]
                p_item.text = item
                p_item.font.size = Pt(22)
                p_item.font.color.rgb = RGBColor(*text_rgb)
                p_item.alignment = PP_ALIGN.CENTER

    def clear_cache(self, older_than_days: int = None):
        """
        Clear cache

        Args:
            older_than_days: Clear cache older than N days, None means clear all
        """
        if self.cache_manager:
            if older_than_days:
                self.cache_manager.cleanup_expired()
                logger.info(f"Cleared expired cache (older than {self.cache_manager.cache_ttl.days} days)")
            else:
                import shutil
                if self.cache_manager.cache_dir.exists():
                    shutil.rmtree(self.cache_manager.cache_dir)
                    logger.info("Cleared all cache")
        else:
            logger.warning("Cache not enabled")

    def get_cache_stats(self) -> Optional[Dict]:
        """
        Get cache statistics

        Returns:
            Dict: Cache statistics
                - outline_count: Outline cache count
                - image_count: Image cache count
                - total_size_mb: Total size (MB)
        """
        if not self.cache_manager:
            return None

        stats = {
            "outline_count": 0,
            "image_count": 0,
            "total_size_mb": 0.0
        }

        try:
            # Count outline cache
            outline_dir = self.cache_manager.outline_cache_dir
            if outline_dir.exists():
                outline_files = list(outline_dir.glob("*.json"))
                stats["outline_count"] = len(outline_files)
                for f in outline_files:
                    stats["total_size_mb"] += f.stat().st_size / (1024 * 1024)

            # Count image cache
            image_dir = self.cache_manager.image_cache_dir
            if image_dir.exists():
                image_files = list(image_dir.glob("*"))
                stats["image_count"] = len(image_files)
                for f in image_files:
                    if f.is_file():
                        stats["total_size_mb"] += f.stat().st_size / (1024 * 1024)

            stats["total_size_mb"] = round(stats["total_size_mb"], 2)

        except Exception as e:
            logger.warning(f"Failed to get cache stats: {e}")

        return stats

    @staticmethod
    def list_template_presets() -> List[Dict]:
        """
        List all available template presets

        Returns:
            List[Dict]: Preset list, each with:
                - key: Preset key (for template_preset parameter)
                - name: Preset name
                - description: Preset description
        """
        return [
            {
                "key": key,
                "name": preset.get("name", key),
                "description": preset.get("description", "")
            }
            for key, preset in get_template_presets().items()
        ]

    @staticmethod
    def get_template_preset_info(preset_name: str) -> Optional[Dict]:
        """
        Get preset details

        Args:
            preset_name: Preset name

        Returns:
            Dict: Preset details with name, description, sequence, narrative
        """
        return get_template_presets().get(preset_name)
